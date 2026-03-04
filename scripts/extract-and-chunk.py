"""
Extract text from downloaded Box files and create search chunks for the ERN Orama index.
Handles: PDF, DOCX, DOC, XLSX, XLSM, PPTX, PPT, TXT
Output: box-chunks.json (same schema as search-chunks.json)
"""
import sys, os, json, re, hashlib
sys.path.insert(0, r'C:\Users\BOB\AppData\Roaming\Python\Python313\site-packages')

DOWNLOADS = r'C:\Users\BOB\.openclaw\workspace\emergency-response\box-downloads'
OUTPUT = r'C:\Users\BOB\.openclaw\workspace\emergency-response\app\src\data\box-chunks.json'
CHUNK_SIZE = 800  # tokens (roughly 4 chars per token)
CHUNK_OVERLAP = 100

stats = {'pdf': 0, 'docx': 0, 'doc': 0, 'xlsx': 0, 'pptx': 0, 'ppt': 0, 'txt': 0, 'skip': 0, 'error': 0}

def extract_pdf(path):
    import pymupdf
    text = []
    try:
        doc = pymupdf.open(path)
        for page in doc:
            text.append(page.get_text())
        doc.close()
    except Exception as e:
        print(f"  ERROR (PDF): {e}")
        stats['error'] += 1
        return ''
    stats['pdf'] += 1
    return '\n'.join(text)

def extract_docx(path):
    from docx import Document
    try:
        doc = Document(path)
        text = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also get table content
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    text.append(' | '.join(cells))
    except Exception as e:
        print(f"  ERROR (DOCX): {e}")
        stats['error'] += 1
        return ''
    stats['docx'] += 1
    return '\n'.join(text)

def extract_xlsx(path):
    import openpyxl
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        text = []
        for ws in wb.worksheets:
            text.append(f"--- Sheet: {ws.title} ---")
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None and str(c).strip()]
                if vals:
                    text.append(' | '.join(vals))
        wb.close()
    except Exception as e:
        print(f"  ERROR (XLSX): {e}")
        stats['error'] += 1
        return ''
    stats['xlsx'] += 1
    return '\n'.join(text)

def extract_pptx(path):
    from pptx import Presentation
    try:
        prs = Presentation(path)
        text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_text.append(' | '.join(cells))
            if slide_text:
                text.append(f"--- Slide {i+1} ---\n" + '\n'.join(slide_text))
    except Exception as e:
        print(f"  ERROR (PPTX): {e}")
        stats['error'] += 1
        return ''
    stats['pptx'] += 1
    return '\n'.join(text)

def extract_txt(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
    except Exception as e:
        print(f"  ERROR (TXT): {e}")
        stats['error'] += 1
        return ''
    stats['txt'] += 1
    return text

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.pdf':
        return extract_pdf(path)
    elif ext in ('.docx', '.doc'):
        if ext == '.doc':
            # Try as docx (some .doc are actually docx)
            try:
                result = extract_docx(path)
                if result:
                    stats['doc'] += 1
                    stats['docx'] -= 1  # undo docx count
                    return result
            except:
                pass
            stats['skip'] += 1
            return ''
        return extract_docx(path)
    elif ext in ('.xlsx', '.xlsm'):
        return extract_xlsx(path)
    elif ext in ('.pptx', '.ppt'):
        if ext == '.ppt':
            stats['skip'] += 1
            return ''  # Can't read old .ppt format
        return extract_pptx(path)
    elif ext == '.txt':
        return extract_txt(path)
    else:
        stats['skip'] += 1
        return ''

def classify_sector(filename, text):
    """Try to guess the sector from filename and content."""
    fn = filename.lower()
    t = text[:2000].lower()
    
    if any(w in fn for w in ['finance', 'budget', 'cam ', 'accounting', 'cost', 'expense', 'journal', 'accounts', 'bank', 'month-end', 'monthend', 'dimension']):
        return 'Finance'
    if any(w in fn for w in ['hr', 'people', 'staff', 'recruit', 'handbook', 'deploy', 'benefit', 'per diem', 'r&r', 'jd', 'job description']):
        return 'People & Culture'
    if any(w in fn for w in ['supply', 'procurement', 'inventory', 'purchase', 'po ', 'item forecast', 'product receipt', 'fixed asset', 'dispatch']):
        return 'Supply Chain'
    if any(w in fn for w in ['safety', 'security', 'sop', 'checkpoint', 'kidnap', 'arrest', 'evacuation', 'fire', 'bomb', 'landmine', 'crossfire', 'carjack', 'robbery', 'grab bag', 'headcount', ' cp ', 'critical incident', 'smp', 'warden', 'proof of life', 'access measurement']):
        return 'Safety & Security'
    if any(w in fn for w in ['safeguard', 'iec ']):
        return 'Safeguarding'
    if any(w in fn for w in ['meal', 'pid', 'measurement', 'indicator', 'assessment report', 'analytical']):
        return 'MEAL'
    if any(w in fn for w in ['partner', 'peers', 'collaborative']):
        return 'Partnerships'
    if any(w in fn for w in ['grant', 'appeal', 'logframe', 'crf', 'donor']):
        return 'Grants'
    if any(w in fn for w in ['sap', 'preparedness', 'early action', 'hazard', 'scenario', 'sitrep', 'crisis response']):
        return 'Response Management'
    if any(w in fn for w in ['integra', 'userguide', 'user guide', 'ug_']):
        return 'Integra'
    if any(w in fn for w in ['quality', 'qie', 'technical assistance']):
        return 'Technical Programs'
    if any(w in fn for w in ['legal', 'counsel']):
        return 'Legal'
    if any(w in fn for w in ['humanitarian access', 'access support']):
        return 'Humanitarian Access'
    
    # Content-based fallback
    if 'finance' in t or 'budget' in t[:500]:
        return 'Finance'
    if 'safety' in t or 'security' in t[:500]:
        return 'Safety & Security'
    
    return 'General'

def chunk_text(text, filename, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    """Split text into overlapping chunks."""
    # Clean up whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    text = text.strip()
    
    if not text or len(text) < 50:
        return []
    
    # Approximate token count (4 chars per token)
    char_size = chunk_size * 4
    char_overlap = overlap * 4
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + char_size
        chunk = text[start:end]
        
        # Try to break at paragraph or sentence boundary
        if end < len(text):
            last_para = chunk.rfind('\n\n')
            if last_para > char_size * 0.5:
                chunk = chunk[:last_para]
                end = start + last_para
            else:
                last_sent = max(chunk.rfind('. '), chunk.rfind('.\n'))
                if last_sent > char_size * 0.5:
                    chunk = chunk[:last_sent + 1]
                    end = start + last_sent + 1
        
        chunk = chunk.strip()
        if chunk and len(chunk) > 30:
            chunks.append(chunk)
        
        start = end - char_overlap if end < len(text) else len(text)
    
    return chunks

# Main
print("Extracting text from downloaded files...")
files = sorted(os.listdir(DOWNLOADS))
all_chunks = []
total_chars = 0
chunk_id = 0

for fname in files:
    fpath = os.path.join(DOWNLOADS, fname)
    if not os.path.isfile(fpath):
        continue
    
    ext = os.path.splitext(fname)[1].lower()
    if ext in ('.zip', '.msg', '.mp4', '.mov'):
        stats['skip'] += 1
        continue
    
    print(f"  {fname}...", end='', flush=True)
    text = extract_text(fpath)
    
    if not text or len(text.strip()) < 50:
        print(f" (empty)")
        continue
    
    total_chars += len(text)
    sector = classify_sector(fname, text)
    chunks = chunk_text(text, fname)
    
    for i, chunk_text_content in enumerate(chunks):
        chunk_id += 1
        cid = f"box-{hashlib.md5(f'{fname}-{i}'.encode()).hexdigest()[:8]}"
        all_chunks.append({
            'id': cid,
            'type': 'document',
            'sector': sector,
            'sectorId': sector.lower().replace(' & ', '-').replace(' ', '-'),
            'phase': 'Reference',
            'title': f"{fname}" + (f" (part {i+1}/{len(chunks)})" if len(chunks) > 1 else ""),
            'content': chunk_text_content,
            'priority': 'reference',
            'source': fname
        })
    
    print(f" {len(chunks)} chunks ({len(text):,} chars) [{sector}]")

# Write output
with open(OUTPUT, 'w', encoding='utf-8') as f:
    json.dump(all_chunks, f, indent=2, ensure_ascii=False)

print(f"\n{'='*60}")
print(f"EXTRACTION COMPLETE")
print(f"  Files processed: {sum(v for k, v in stats.items() if k != 'skip' and k != 'error')}")
print(f"  Skipped: {stats['skip']}, Errors: {stats['error']}")
print(f"  Total text: {total_chars:,} chars (~{total_chars//4:,} tokens)")
print(f"  Chunks created: {len(all_chunks)}")
print(f"  Output: {OUTPUT}")
print(f"\nBy format: {json.dumps(stats)}")
